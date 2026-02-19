from moviepy import VideoFileClip
import os
import streamlit as st
# ENABLE WIDE MODE (Must be first Streamlit command)
st.set_page_config(page_title="Viral Engine v2", layout="wide", initial_sidebar_state="expanded")

import instaloader
import yt_dlp
import base64
import assemblyai as aai
from dotenv import load_dotenv
import json
import io
import requests
import google.generativeai as genai
from PIL import Image
try:
    import whisper
except ImportError:
    st.error("OpenAI Whisper not installed. Please run `pip install openai-whisper`")
    whisper = None # Handle missing whisper gracefully

from database import DatabaseManager

# Load environment variables
env_path = "/Users/jaydengle/Transcribe-Reels/.env"
load_dotenv(dotenv_path=env_path)

# Ensure external tools (ffmpeg) are found
os.environ["PATH"] += os.pathsep + "/opt/homebrew/bin" + os.pathsep + "/usr/local/bin"

# Set AssemblyAI API key
aai.settings.api_key = os.getenv("ASSEMBLYAI_API_KEY")

# Set Google Gemini API key
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Initialize Database
db_manager = DatabaseManager()
db_connected = db_manager.connect()

if not db_connected:
    st.warning("Could not connect to Supabase. Transcripts will not be saved.")

# Function to download Instagram Reel & Extract Metadata
def download_reel(url):
    # Try yt-dlp first with HARDENED HEADERS
    try:
        if not os.path.exists('reel'): os.makedirs('reel')
        for f in os.listdir('reel'): 
            try: os.remove(os.path.join('reel', f))
            except: pass

        ydl_opts = {
            'format': 'bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best',
            'outtmpl': 'reel/%(title)s.%(ext)s',
            'quiet': True, 'no_warnings': True,
            'user_agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'referer': 'https://www.instagram.com/',
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            # Find any video file (mp4, mkv, webm)
            video_exts = ['.mp4', '.mkv', '.webm', '.mov']
            downloaded = [f for f in os.listdir('reel') if any(f.endswith(ext) for ext in video_exts)]
            if downloaded:
                return {
                    "title": info.get('title', 'video'),
                    "caption": info.get('description', ''),
                    "owner": info.get('uploader', 'Unknown'),
                    "video_path": os.path.join('reel', downloaded[0])
                }

    except Exception as e:
        # Silently fallback without scary errors
        pass


    # Fallback to Instaloader for IG specifically
    try:
        loader = instaloader.Instaloader()
        if 'instagram.com/reel/' in url:
            shortcode = url.split('/')[-2]
            post = instaloader.Post.from_shortcode(loader.context, shortcode)
            loader.download_post(post, target='reel')
            downloaded = [f for f in os.listdir('reel') if f.endswith('.mp4')]
            if downloaded:
                return {
                    "title": f"Reel {shortcode}",
                    "caption": post.caption or "",
                    "owner": post.owner_username,
                    "video_path": os.path.join('reel', downloaded[0])
                }
    except Exception as e:
        st.error(f"Download totally failed: {e}")
    return None

# Function to analyze video visuals using GPT-4o (Vision)
import base64
def analyze_visual_frames(video_path):
    clip = VideoFileClip(video_path)
    duration = clip.duration
    
    # Extract frames at critical points (Hook, Mid, End)
    timestamps = [0, min(3, duration), min(10, duration), min(duration-2, duration)]
    unique_timestamps = sorted(list(set(timestamps)))
    
    frames_base64 = []
    for t in unique_timestamps:
        frame_path = f"temp_frame_{t}.jpg"
        clip.save_frame(frame_path, t=t)
        
        with open(frame_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            frames_base64.append(encoded_string)
        
        if os.path.exists(frame_path):
            os.remove(frame_path) # Cleanup
            
    clip.close()
    return frames_base64

# Function to convert video to audio
def convert_video_to_audio(video_file_path):
    clip = VideoFileClip(video_file_path)
    clip.audio.write_audiofile("audio.wav", codec='pcm_s16le', ffmpeg_params=["-ac", "1", "-ar", "16000"]) 
    clip.close()
    return "audio.wav"

# Function to transcribe utilizing Whisper (Local)
def transcribe_with_whisper(audio_file, model_size="base"):
    if whisper is None:
        st.error("Whisper library not found.")
        return None
    
    model = whisper.load_model(model_size)
    # Balanced Settings (Allowing more context but with conservative speech detection)
    result = model.transcribe(audio_file, condition_on_previous_text=True, no_speech_threshold=0.6)
    return result

# --- 🧠 GEMINI SKILL EXTRACTION ---
def extract_skills_gemini(frames):
    """Extract technical skills/actions from video frames using Gemini Flash Latest."""
    try:
        # Use gemini-flash-latest for stable access
        model = genai.GenerativeModel('gemini-flash-latest')
        
        # Convert base64 frames to PIL Images for Gemini
        images = []
        for f in frames[:10]: # Limit to 10 frames for cost/speed
            img_data = base64.b64decode(f)
            images.append(Image.open(io.BytesIO(img_data)))
        
        prompt = """
        You are a Technical Skill Analyst. Analyze this sequence of frames from a video.
        
        Provide a detailed breakdown of HOW they did what they did in each part of the video:
        1. **Phase/Part**: Identify the distinct stages or segments of the action.
        2. **Technical Details**: For each stage, describe the exact body positioning, tool usage, or specific technique applied.
        3. **Hidden 'Pro' Cues**: Point out subtle details that make the skill effective.
        
        Format the output clearly as a step-by-step technical guide.
        """
        
        response = model.generate_content([prompt] + images)
        return response.text
    except Exception as e:
        return f"Gemini Extraction failed: {e}"

# --- 🎨 GEMINI IMAGE GENERATION (PROMPT ENGINE) ---
def generate_image_desc_gemini(analysis):
    """Uses Gemini to create a highly detailed prompt for an image that represents the 'vibe' or 'skill'."""
    try:
        model = genai.GenerativeModel('gemini-flash-latest')
        prompt = f"Based on this skill analysis: '{analysis}', describe a premium, high-resolution aesthetic image that represents this skill. Focus on lighting, composition, and mood. Output ONLY the description."
        response = model.generate_content(prompt)
        return response.text
    except:
        return "A professional aesthetic representation of technical skills."

# --- 🎨 GEMINI IMAGE GENERATION ---
def generate_image_pollinations(prompt):
    """Generates an image URL using Pollinations.ai (Free, Unlimited, No Key)."""
    try:
        # Encode prompt for URL
        import urllib.parse
        encoded_prompt = urllib.parse.quote(prompt)
        # Enhance prompt for better results
        enhanced_prompt = f"{encoded_prompt} high quality, detailed, 8k, cinematic lighting"
        
        # Pollinations.ai URL format
        image_url = f"https://image.pollinations.ai/prompt/{enhanced_prompt}?width=1024&height=1024&nologo=true&model=flux"
        return image_url
    except Exception as e:
        st.error(f"Pollinations Image Gen failed: {e}")
        return None

# --- 🎙️ EDUCATIONAL PODCAST GENERATOR (NotebookLM-Style) ---
def generate_educational_podcast(transcript_text, analysis_text, openai_client):
    """Generates a 3-minute dual-speaker educational podcast using GPT + OpenAI TTS."""
    try:
        from pydub import AudioSegment
        import tempfile

        # Step 1: Generate the dialogue script (topic + skills focus)
        script_prompt = f"""
You are a podcast script writer for a super chill, conversational deep-dive podcast — think two friends who really know their stuff just riffing and teaching each other.

Write a ~3-minute casual educational conversation (about 600 words) between:
- ALEX: The host. Curious, laid-back, asks great questions. Uses casual language like "wait...", "oh that's actually sick", "right?", "so basically...", "okay okay"
- JAMIE: The expert. Naturally enthusiastic about the topic. Breaks things down step-by-step but keeps it real and relatable. Uses phrases like "yeah exactly", "so the thing is...", "what most people don't realize is...", "honestly..."

IMPORTANT RULES:
- Write like real people talking, NOT like a formal presentation. Use contractions, natural pauses (...), and filler words.
- The podcast is ABOUT THE TOPIC AND SKILLS in the video — NOT about why the video went viral.
- JAMIE should deeply expand on the specific techniques and skills shown. Give real detail.
- Keep sentences short and punchy. No big walls of text per turn.
- Use natural turn-taking — it should feel like a flowing conversation.

Structure:
1. ALEX casually introduces the topic of the video (not the video itself)
2. JAMIE gives a warm overview of what the skill/topic is about
3. ALEX asks about the most interesting or surprising technique
4. JAMIE breaks down the key steps in a natural, excited way
5. Back-and-forth on details, tips, and things most people get wrong
6. ALEX wraps up with a "so what's the one thing someone should take away from this?"

Format EXACTLY like this:
ALEX: [line]
JAMIE: [line]
ALEX: [line]
...

VIDEO TRANSCRIPT: {transcript_text[:500]}
SKILL/TOPIC CONTEXT: {analysis_text[:800]}
"""
        script_resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": script_prompt}],
            max_tokens=1100
        )
        script = script_resp.choices[0].message.content
        # Normalize labels in case GPT used the new names
        script = script.replace("ALEX:", "HOST:").replace("JAMIE:", "EXPERT:")

        # Step 2: Parse lines and generate TTS per speaker
        lines = [l.strip() for l in script.split('\n') if l.strip() and ':' in l]
        base_dir = "/Users/jaydengle/Transcribe-Reels"
        clips = []
        silence = AudioSegment.silent(duration=500)  # natural pause between turns

        for i, line in enumerate(lines):
            if line.startswith("HOST:"):
                voice = "nova"   # warm, natural female voice
                text = line[5:].strip()
            elif line.startswith("EXPERT:"):
                voice = "onyx"   # deep, confident male voice
                text = line[7:].strip()
            else:
                continue
            if not text:
                continue

            tts_resp = openai_client.audio.speech.create(
                model="tts-1-hd",  # higher quality, more natural
                voice=voice,
                input=text
            )
            tmp_path = os.path.join(base_dir, f"_tmp_clip_{i}.mp3")
            tts_resp.stream_to_file(tmp_path)
            clips.append(AudioSegment.from_mp3(tmp_path) + silence)

        if not clips:
            return None

        # Step 3: Stitch and save
        podcast = clips[0]
        for c in clips[1:]:
            podcast += c

        out_path = os.path.join(base_dir, "educational_podcast.mp3")
        podcast.export(out_path, format="mp3")

        # Cleanup temp clips
        for i in range(len(lines)):
            p = os.path.join(base_dir, f"_tmp_clip_{i}.mp3")
            if os.path.exists(p): os.remove(p)

        return out_path
    except Exception as e:
        return f"ERROR: {e}"
# --- 🔍 REVERSE ENGINEER ANY SCREENSHOT ---
def reverse_engineer_image(image_bytes, openai_client):
    """Sends an image to GPT-4o Vision to identify the tool/technique and provide replication steps."""
    try:
        b64 = base64.b64encode(image_bytes).decode('utf-8')
        prompt_content = [
            {
                "type": "text",
                "text": """You are an expert AI tools analyst. Look at this screenshot carefully.

Provide a detailed, structured breakdown in this exact format:

## 🔧 Tool / Software Identified
Name the exact tool, platform, or software shown. Be specific.

## 📺 What's Happening
Describe exactly what workflow, technique, or feature is being used.

## 📄 How To Replicate This (Step-by-Step)
1. Step one...
2. Step two...
(give detailed, practical steps a beginner can follow)

## 🛠️ Tools & Software You Need
- Tool name: description + link if known

## 💰 Estimated Cost
Free / Paid / Freemium — specify pricing tier if known.

## 🔍 GitHub Search Terms
Provide 2-3 search terms to find related repos on GitHub (e.g. 'kling-ai comfyui workflow')."""
            },
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}}
        ]
        resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt_content}],
            max_tokens=800
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"ERROR: {e}"

# --- 📦 GITHUB SKILL DISCOVERY ---
def search_github_for_skill(query):
    """Searches GitHub for repos related to a tool/skill and returns the top results."""
    try:
        url = f"https://api.github.com/search/repositories?q={requests.utils.quote(query)}&sort=stars&order=desc&per_page=6"
        headers = {"Accept": "application/vnd.github.v3+json", "User-Agent": "ViralEngine/1.0"}
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code != 200:
            return []
        items = resp.json().get('items', [])
        return [{
            "name": r["full_name"],
            "description": r.get("description", "No description"),
            "stars": r["stargazers_count"],
            "url": r["html_url"],
            "clone": r["clone_url"]
        } for r in items]
    except Exception as e:
        return []

# --- 🎥 KLING AI STORYBOARD GENERATOR (Scene Detection + Dense 2s Sampling) ---
def generate_kling_storyboard(video_path, openai_client):
    """Analyzes video using PySceneDetect (real shot boundaries) + dense 2s sampling for maximum accuracy."""
    try:
        from moviepy import VideoFileClip
        import cv2

        base_dir = "/Users/jaydengle/Transcribe-Reels"
        clip = VideoFileClip(video_path)
        duration = clip.duration

        # ── METHOD 1: PySceneDetect — find real shot/cut boundaries ──
        scene_timestamps = []
        try:
            # Lazy import to avoid startup crashes if missing
            import scenedetect
            from scenedetect import open_video, SceneManager
            from scenedetect.detectors import ContentDetector
            
            video = open_video(video_path)
            scene_mgr = SceneManager()
            scene_mgr.add_detector(ContentDetector(threshold=27.0))
            scene_mgr.detect_scenes(video, show_progress=False)
            scene_list = scene_mgr.get_scene_list()
            scene_timestamps = [s[0].get_seconds() for s in scene_list]
        except Exception as sd_e:
            st.warning(f"⚠️ Scene detection skipped (using time sampling): {sd_e}")
            scene_timestamps = []  # fallback to dense sampling below

        # ── METHOD 2: Dense 2-second sampling (catches subtle changes) ──
        dense_timestamps = list(range(0, int(duration), 2))

        # ── MERGE: combine both, deduplicate, sort ──
        all_ts = sorted(set(int(t) for t in scene_timestamps + dense_timestamps))
        # Cap at 24 frames to keep cost manageable
        if len(all_ts) > 24:
            # Smart downsample: keep scene boundaries + spread the rest
            scene_set = set(int(t) for t in scene_timestamps)
            kept = list(scene_set)
            remaining = [t for t in all_ts if t not in scene_set]
            step = max(1, len(remaining) // (24 - len(kept)))
            kept += remaining[::step]
            all_ts = sorted(set(kept))[:24]

        # ── EXTRACT FRAMES ──
        shots = []
        for t in all_ts:
            if t >= duration: continue
            frame_path = os.path.join(base_dir, f"_kling_frame_{t}.jpg")
            clip.save_frame(frame_path, t=t)
            with open(frame_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode('utf-8')
            shots.append({"t": t, "path": frame_path, "b64": b64})
        clip.close()

        # ── GPT-4o VISION: write Kling prompt per shot ──
        storyboard = []
        for i, shot in enumerate(shots):
            next_t = shots[i+1]["t"] if i+1 < len(shots) else int(duration)
            duration_s = max(1, next_t - shot["t"])
            is_cut = shot["t"] in [int(t) for t in scene_timestamps]

            prompt_content = [
                {
                    "type": "text",
                    "text": f"""You are a professional AI video director writing Kling AI Story Mode prompts.

This is shot {i+1} captured at {shot['t']}s ({f'⚡ SCENE CUT DETECTED' if is_cut else 'continuous shot'}).
Estimated shot duration: ~{min(duration_s, 10)}s.

Study the image VERY carefully. Describe ONLY what you actually see — no guessing.

Write ONE Kling AI cinematic prompt using this exact structure:
[Subject doing X], [specific setting/environment], [camera move], [lighting style], [mood], cinematic, photorealistic, 4K, smooth motion, no camera shake. Duration: {min(duration_s, 10)}s.

Film language to use:
- Camera: 'slow dolly in', 'push in', 'wide static shot', 'handheld tracking', 'aerial pull back', 'rack focus'
- Lighting: 'harsh rim light', 'golden hour', 'studio soft box', 'neon glow', 'natural daylight'

Output ONLY the prompt text. No labels, no explanation."""
                },
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{shot['b64']}"}}
            ]
            resp = openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": prompt_content}],
                max_tokens=220
            )
            kling_prompt = resp.choices[0].message.content.strip()
            storyboard.append({
                "shot": i + 1,
                "timestamp": shot["t"],
                "duration": min(duration_s, 10),
                "prompt": kling_prompt,
                "is_scene_cut": is_cut
            })

        # Cleanup temp frames
        for shot in shots:
            if os.path.exists(shot["path"]): os.remove(shot["path"])

        return storyboard
    except Exception as e:
        return f"ERROR: {e}"

# --- 🎥 KLING AI STORYBOARD GENERATOR (Gemini Native) ---
def generate_kling_storyboard_gemini(video_path):
    """
    Uses Gemini 1.5 Flash (Native Video) for free, fast, unlimited storyboard generation.
    Bypasses local frame extraction and GPT-4o costs.
    """
    try:
        import time
        
        # 1. Upload Video to Gemini File API
        print(f"Uploading {video_path} to Gemini...")
        video_file = genai.upload_file(path=video_path)
        
        # 2. Wait for processing
        while video_file.state.name == "PROCESSING":
            time.sleep(2)
            video_file = genai.get_file(video_file.name)
            
        if video_file.state.name == "FAILED":
            raise ValueError("Gemini Video Processing Failed")
            
        # 3. Prompt for JSON Storyboard
        # Use gemini-flash-latest alias for best availability
        model = genai.GenerativeModel('gemini-flash-latest')
        
        prompt = """
        You are a professional Film Director. Analyze this video shot-by-shot.
        Return a JSON list of objects, where each object represents a distinct shot/scene.
        
        JSON Format:
        [
            {
                "shot": 1,
                "timestamp": 0,
                "duration": 4,
                "prompt": "[Subject doing X], [Environment], [Type of Shot], [Lighting], cinematic, 8k"
            },
            ...
        ]
        
        Rules for 'prompt':
        - MUST be a high-quality prompt for Kling AI / Sora.
        - Include Camera Movement (e.g., "Slow dolly in", "Truck left").
        - Include Lighting (e.g., "Golden hour", "Cyberpunk neon").
        - Describe EXACTLY what is visible.
        - No text overlays.
        """
        
        # Request JSON mode
        response = model.generate_content([video_file, prompt], generation_config={"response_mime_type": "application/json"})
        
        # Parse JSON
        import json
        storyboard_data = json.loads(response.text)
        return storyboard_data if isinstance(storyboard_data, list) else []
        
    except Exception as e:
        st.error(f"Gemini Storyboard Error: {e}")
        return str(e)

# --- 📚 VIDEO-TO-PDF HOW-TO GUIDE ---
def generate_how_to_pdf(video_path, transcript_text, openai_client):
    """Reads each video frame section with GPT-4o Vision and builds a step-by-step PDF how-to guide."""
    try:
        from fpdf import FPDF
        from moviepy import VideoFileClip

        base_dir = "/Users/jaydengle/Transcribe-Reels"
        clip = VideoFileClip(video_path)
        duration = clip.duration
        
        # Extract frames every 4 seconds (max 20 sections)
        interval = max(4, int(duration / 15))
        frame_data = []
        
        for t in range(0, int(duration), interval):
            if len(frame_data) >= 18: break
            frame_path = os.path.join(base_dir, f"_pdf_frame_{t}.jpg")
            clip.save_frame(frame_path, t=t)
            with open(frame_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode('utf-8')
            frame_data.append({"t": t, "path": frame_path, "b64": b64})
        clip.close()

        # Send all frames to GPT-4o Vision for per-section analysis
        step_descriptions = []
        for i, fd in enumerate(frame_data):
            prompt_content = [
                {"type": "text", "text": f"You are creating a step-by-step how-to guide. This is frame {i+1} at timestamp {fd['t']}s of the video.\n\nDescribe EXACTLY what is happening here as a practical guide step. Be specific about:\n- What action is being performed\n- Body position, technique, or tool being used\n- Key detail a beginner would need to know\n\nFormat: Start with 'Step {i+1}:' then 2-3 sentences max.\n\nVideo transcript context: {transcript_text[:300]}"},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{fd['b64']}"}}
            ]
            resp = openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": prompt_content}],
                max_tokens=150
            )
            step_descriptions.append(resp.choices[0].message.content.strip())

        # Build the PDF
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        
        # Title
        pdf.set_font("Helvetica", "B", 22)
        pdf.set_text_color(30, 30, 30)
        title = f"How-To Guide"
        pdf.cell(0, 12, title, ln=True, align="C")
        pdf.set_font("Helvetica", "", 11)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 8, f"Generated from video | {len(frame_data)} sections analyzed by GPT-4o Vision", ln=True, align="C")
        pdf.ln(4)
        pdf.set_draw_color(200, 0, 0)
        pdf.set_line_width(0.5)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(6)

        # Transcript Summary
        if transcript_text:
            pdf.set_font("Helvetica", "B", 13)
            pdf.set_text_color(30, 30, 30)
            pdf.cell(0, 8, "What This Video Is About", ln=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(60, 60, 60)
            pdf.multi_cell(0, 6, transcript_text[:400])
            pdf.ln(4)

        # Step-by-step
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(180, 0, 0)
        pdf.cell(0, 10, "Step-by-Step Guide", ln=True)
        pdf.ln(2)

        for i, (fd, desc) in enumerate(zip(frame_data, step_descriptions)):
            # Check page space
            if pdf.get_y() > 240:
                pdf.add_page()
            
            # Layout: image left, text right
            img_x, img_y = 10, pdf.get_y()
            img_w, img_h = 60, 45
            
            # Embed frame
            try:
                pdf.image(fd["path"], x=img_x, y=img_y, w=img_w, h=img_h)
            except: pass

            # Step text
            pdf.set_xy(img_x + img_w + 5, img_y)
            pdf.set_font("Helvetica", "B", 11)
            pdf.set_text_color(180, 0, 0)
            pdf.cell(0, 7, f"Section {i+1}  [{fd['t']}s]", ln=True)
            pdf.set_x(img_x + img_w + 5)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(40, 40, 40)
            pdf.multi_cell(120, 6, desc)
            
            # Advance past the image
            pdf.set_y(max(pdf.get_y(), img_y + img_h + 4))
            pdf.ln(2)

        # Tips footer
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(180, 0, 0)
        pdf.cell(0, 10, "Pro Tips & Summary", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(50, 50, 50)
        tips_prompt = f"Based on this video transcript, write 5 concise pro tips a beginner should know. Transcript: {transcript_text[:600]}"
        tips_resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": tips_prompt}],
            max_tokens=300
        )
        pdf.multi_cell(0, 7, tips_resp.choices[0].message.content)

        # Save PDF
        out_path = os.path.join(base_dir, "how_to_guide.pdf")
        pdf.output(out_path)

        # Cleanup frames
        for fd in frame_data:
            if os.path.exists(fd["path"]): os.remove(fd["path"])

        return out_path
    except Exception as e:
        return f"ERROR: {e}"

# Function to transcribe utilizing AssemblyAI (API)
def transcribe_with_assemblyai(audio_file):
    transcriber = aai.Transcriber()
    config = aai.TranscriptionConfig(
        speaker_labels=True,
        speech_models=["universal-2"],
        language_detection=True
    )
    with open(audio_file, "rb") as f:
        transcript = transcriber.transcribe(f, config)
    return transcript

# NEW: Detailed Visual Timeline Extractor
def generate_visual_timeline(video_path, interval=3):
    clip = VideoFileClip(video_path)
    duration = clip.duration
    frames_base64 = []
    timestamps = []
    
    # Extract frame every 'interval' seconds
    for t in range(0, int(duration), interval):
        if t > 60: break # Limit to first 60s to save tokens/time for now
        
        frame_path = f"temp_timeline_{t}.jpg"
        clip.save_frame(frame_path, t=t)
        
        with open(frame_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            frames_base64.append(encoded_string)
            timestamps.append(t)
        
        if os.path.exists(frame_path): os.remove(frame_path)
            
    clip.close()
    return frames_base64, timestamps

# Initialize OpenAI

# Initialize OpenAI
openai_api_key = os.getenv("OPENAI_API_KEY")
if openai_api_key:
    from openai import OpenAI
    openai_client = OpenAI(api_key=openai_api_key)
else:
    openai_client = None

# Custom CSS for "Jay's Cheat Model" Theme (Red/Black Hacker Vibe)
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Space+Mono:ital,wght@0,400;0,700;1,400&display=swap');

    /* Main Background - Dark Cyber */
    .stApp {
        background-color: #000000;
        background-image: radial-gradient(circle at 50% 50%, #1a0000 0%, #000000 100%);
        color: #ffcccc;
        font-family: 'Space Mono', monospace;
    }
    
    /* Sidebar Styling */
    [data-testid="stSidebar"] {
        background-color: #050000;
        border-right: 1px solid #330000;
    }
    
    /* Card/Container Glassmorphism - Red Hint */
    div.css-1r6slb0.e1tzin5v2, [data-testid="stExpander"] {
        background: rgba(20, 0, 0, 0.7);
        backdrop-filter: blur(5px);
        border: 1px solid #ff0000;
        border-radius: 4px; /* Sharper edges for tech feel */
        box-shadow: 0 0 15px rgba(255, 0, 0, 0.1);
    }
    
    /* Glitch Title Style */
    .glitch-title {
        font-size: 3rem;
        font-weight: bold;
        text-transform: uppercase;
        position: relative;
        text-shadow: 0.05em 0 0 #00fffc, -0.03em -0.04em 0 #fc00ff,
                     0.025em 0.04em 0 #fffc00;
        animation: glitch 725ms infinite;
        color: #ffffff;
    }
    
    @keyframes glitch {
        0% { text-shadow: 0.05em 0 0 #00fffc, -0.03em -0.04em 0 #fc00ff, 0.025em 0.04em 0 #fffc00; }
        15% { text-shadow: 0.05em 0 0 #00fffc, -0.03em -0.04em 0 #fc00ff, 0.025em 0.04em 0 #fffc00; }
        16% { text-shadow: -0.05em -0.025em 0 #00fffc, 0.025em 0.035em 0 #fc00ff, -0.05em -0.05em 0 #fffc00; }
        49% { text-shadow: -0.05em -0.025em 0 #00fffc, 0.025em 0.035em 0 #fc00ff, -0.05em -0.05em 0 #fffc00; }
        50% { text-shadow: 0.05em 0.035em 0 #00fffc, 0.03em 0 0 #fc00ff, 0 -0.04em 0 #fffc00; }
        99% { text-shadow: 0.05em 0.035em 0 #00fffc, 0.03em 0 0 #fc00ff, 0 -0.04em 0 #fffc00; }
        100% { text-shadow: -0.05em 0 0 #00fffc, -0.025em -0.04em 0 #fc00ff, -0.04em -0.025em 0 #fffc00; }
    }
    
    /* Headers - MAX VISIBILITY RED */
    h1, h2, h3, h4, h5, h6 {
        color: #ff0000 !important;
        text-transform: uppercase;
        font-family: 'Space Mono', monospace !important;
        letter-spacing: 1px;
        text-shadow: 0 0 10px rgba(255, 0, 0, 0.4);
    }
    
    /* Make bold text (like "1. Hook:") stand out in Red */
    strong {
        color: #ff3333 !important;
        font-weight: 900;
    }

    /* Buttons - Red Neon */
    .stButton > button {
        background: #000000;
        color: #ff0000;
        border: 1px solid #ff0000;
        border-radius: 0px; /* Cyberpunk style */
        padding: 0.5rem 1rem;
        font-weight: bold;
        text-transform: uppercase;
        transition: all 0.2s ease;
        box-shadow: 0 0 10px rgba(255, 0, 0, 0.2);
    }
    .stButton > button:hover {
        background: #ff0000;
        color: #000000;
        box-shadow: 0 0 20px rgba(255, 0, 0, 0.8);
    }
    
    /* Input Fields */
    .stTextInput > div > div > input {
        background-color: #0a0000;
        color: #ff0000;
        border: 1px solid #330000;
        border-radius: 0px;
        font-family: 'Space Mono', monospace;
    }
    .stTextInput > div > div > input:focus {
        border-color: #ff0000;
        box-shadow: 0 0 10px rgba(255, 0, 0, 0.4);
    }
    
    /* Expander Transitions */
    .streamlit-expanderHeader {
        background-color: transparent;
        color: #ffcccc;
        font-family: 'Space Mono', monospace;
    }
</style>
""", unsafe_allow_html=True)

# Main UI Structure
# Use custom HTML for the Glitch Title Logo
st.markdown('<h1 class="glitch-title">JAY\'S VIRAL ENGINE ⚡</h1>', unsafe_allow_html=True)

# Sidebar for controls (restoring missing settings)
st.sidebar.header("SYSTEM CONFIG")
transcription_engine = st.sidebar.selectbox("ENGINE CORE", ["OpenAI Whisper (Local - Better Quality)", "AssemblyAI (Cloud)"])
model_size = "base"
if transcription_engine == "OpenAI Whisper (Local - Better Quality)":
    model_size = st.sidebar.selectbox("MODEL SIZE", ["tiny", "base", "small", "medium", "large"], index=2)

# --- LAYOUT: 2-COLUMN WIDE (Adjusted for better spacing) ---
control_col, workspace_col = st.columns([1.2, 3])

# --- LEFT: CONTROL CENTER (Sources & Skills) ---
with control_col:
    # HELP GUIDE:
    with st.expander("🗺️ USER GUIDE (Start Here)"):
        st.markdown("""
        **1. ⚡ New Analysis**: Enter URL below -> Click 'ANALYZE'.
        **2. 🕸️ Mind Map**: Go to **Tab 2 (Viral DNA)** -> 'Visual Mind Map'.
        **3. 🖼️ Cover Art**: Go to **Tab 4 (Visuals)** -> 'Generate Cover Art'.
        **4. 📊 Slide Deck**: Go to **Tab 4 (Visuals)** -> 'Create Slide Deck'.
        **5. 💬 AI Chat**: Go to **Tab 6 (AI Chat)**.
        **6. 🔊 Audio**: Check Player below if transcript fails.
        """)
        
    st.markdown("### 🗂️ LIBRARY")
    
    # --- 1. NEW IMPORT WORKFLOW (Open Layout) ---
    st.markdown("#### ➕ New Source")
    url = st.text_input("URL", placeholder="Instagram/YouTube Link...", label_visibility="collapsed")
    
    # Step 1: Analyze & Preview
    if st.button("⚡ ANALYZE", use_container_width=True):
        if url:
            with st.spinner("🚀 Downloading & extracting audio..."):
                # Clean previous session preview
                if 'preview_reel' in st.session_state: del st.session_state['preview_reel']
                
                reel_data = download_reel(url) # Returns dict or None
                if reel_data:
                    video_path = reel_data["video_path"]
                    with st.spinner("🧠 Transcribing & Categorizing..."):
                        audio_file = convert_video_to_audio(video_path)
                        
                        # audio_file already generated
                        
                        # Transcribe
                        if transcription_engine == "OpenAI Whisper (Local - Better Quality)":
                            result = transcribe_with_whisper(audio_file, model_size)
                            transcript_text = result["text"] if result else ""
                        else:
                            result = transcribe_with_assemblyai(audio_file)
                            transcript_text = result.text if result else ""

                        # Detect Potential Hallucinations (don't delete, just warn)
                        JUNK_TERMS = [
                            "SNOWY", "THANK YOU", "SHOWER", "SUBSCRIBE", "THANKS FOR WATCHING",
                            "PAST IN THE GAMMA'S", "ANYWAY IM THE FAMOUS", "THE END", "THANKS FOR",
                            "THANK YOU SO MUCH", "BE SURE TO LIKE", "TICK TOK"
                        ]
                        
                        st.session_state['hallucination_risk'] = any(term in transcript_text.strip().upper() for term in JUNK_TERMS) or len(transcript_text.strip()) < 15

                        if transcript_text or video_path: # Run analysis even if no transcript but we have video
                            # AI Categorization & Auto-Title

                            auto_title = "Untitled Reel"
                            auto_category = "General"
                            if openai_client:
                                try:
                                    sys_p = "You are a Librarian."
                                    usr_p = f"Suggest a Title (max 5 words) and a Category (1 word) for this content:\n{transcript_text[:500]}"
                                    resp = openai_client.chat.completions.create(model="gpt-3.5-turbo", messages=[{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}])
                                    ai_suggestion = resp.choices[0].message.content
                                    # Simple parsing if format is "Title: X\nCategory: Y"
                                    parts = ai_suggestion.split('\n')
                                    for p in parts:
                                        if "Title:" in p: auto_title = p.replace("Title:", "").strip()
                                        if "Category:" in p: auto_category = p.replace("Category:", "").strip()
                                except: pass

                            # --- 🚀 FULL AUTO-ANALYSIS FOR PREVIEW ---
                            analysis_result = ""
                            if openai_client:
                                try:
                                    # 1. AUTO-VISUALS (Timeline & Storyboard)
                                    status_visual = st.empty()
                                    status_visual.info("🎞️ Generating Visual Timeline...")
                                    frames, stamps = generate_visual_timeline(video_path, interval=5)
                                    
                                    # Quick Visual Summary for Analysis Context
                                    prompt_content = [{"type": "text", "text": "Describe the visual action in these frames shot-by-shot for a viral analyst."}]
                                    for i, f in enumerate(frames[:6]): # Limit for prompt efficiency
                                        prompt_content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{f}"}})
                                    
                                    vis_resp = openai_client.chat.completions.create(model="gpt-4o", messages=[{"role": "user", "content": prompt_content}], max_tokens=400)
                                    visual_analysis_text = vis_resp.choices[0].message.content
                                    
                                    st.session_state['visual_timeline'] = {
                                        "analysis": visual_analysis_text,
                                        "frames": frames, "timestamps": stamps
                                    }
                                    status_visual.success("🎞️ Visual Timeline Ready!")

                                    # 2. VIRAL DNA ANALYSIS
                                    status_dna = st.empty()
                                    status_dna.info("🧬 Analyzing Viral DNA...")
                                    sys_prompt = "You are a Viral Content Analyst."
                                    usr_prompt = f"""
                                    Analyze this content (Transcript + Visual Actions) and output a report with these Exact Headers:

                                    ## 🪝 VIRAL HOOKS
                                    (List 3 alternative hooks)

                                    ## 🧠 WHY IT WENT VIRAL
                                    (Psychological triggers)

                                    ## 📐 CONTENT STRATEGY
                                    (Hook -> Story -> CTA)

                                    TRANSCRIPT: {transcript_text}
                                    VISUAL ACTION: {visual_analysis_text}
                                    """
                                    resp = openai_client.chat.completions.create(model="gpt-4", messages=[{"role": "system", "content": sys_prompt}, {"role": "user", "content": usr_prompt}])
                                    analysis_result = resp.choices[0].message.content
                                    st.session_state['viral_analysis'] = analysis_result
                                    status_dna.success("🧬 Viral DNA Encoded!")

                                    # 3. AUTO MIND MAP
                                    status_map = st.empty()
                                    status_map.info("🕸️ Mapping Flow...")
                                    dot_prompt = f"Create a Graphviz DOT code flow (Hook, Context, Problem, Solution, CTA) based on this analysis:\n{analysis_result}\nOutput DOT only. Begin with 'digraph'."
                                    m_resp = openai_client.chat.completions.create(model="gpt-4o", messages=[{"role": "user", "content": dot_prompt}])
                                    dot_code = m_resp.choices[0].message.content.replace("```dot", "").replace("```", "").replace("dot", "").strip()
                                    st.session_state['viral_map'] = dot_code
                                    status_map.success("🕸️ Flow Map Generated!")

                                    # 4. TTS AUDIO PODCAST
                                    status_tts = st.empty()
                                    status_tts.info("🎙️ Generating Audio Podcast...")
                                    try:
                                        tts_script = f"Here is your viral breakdown. {analysis_result[:1000]}"
                                        tts_response = openai_client.audio.speech.create(
                                            model="tts-1",
                                            voice="onyx",
                                            input=tts_script
                                        )
                                        tts_out_path = "/Users/jaydengle/Transcribe-Reels/overview_audio.mp3"
                                        tts_response.stream_to_file(tts_out_path)
                                        st.session_state['tts_audio_bytes'] = open(tts_out_path, 'rb').read()
                                        status_tts.success("🎙️ Audio Podcast Ready!")
                                    except Exception as tts_e:
                                        status_tts.warning(f"TTS unavailable: {tts_e}")

                                except Exception as e:
                                    st.warning(f"Feature Analysis failed: {e}")
                                    analysis_result = f"Analysis limited: {e}"


                            # Store in Session State for Review
                            st.session_state['preview_source'] = {
                                'url': url,
                                'text': transcript_text,
                                'metadata': reel_data,
                                'title': auto_title,
                                'category': auto_category,
                                'summary': analysis_result # Store analysis here
                            }
                            # FORCE POPULATE MAIN UI STATE
                            st.session_state['viral_analysis'] = analysis_result
                            
                            st.success("Ready to Review!")
                            st.rerun()
                else:
                    st.error("Download failed.")

        # Step 2: Review & Save
        if 'preview_source' in st.session_state:
            p_source = st.session_state['preview_source']
            st.divider()
            st.caption("📝 REVIEW BEFORE SAVING")
            st.text_area("Transcript Preview", value=p_source['text'][:200]+"...", height=100, disabled=True)
            
            p_title = st.text_input("Title", value=p_source.get('title', ''))
            p_project = st.text_input("Project", value="My First Project")
            p_category = st.text_input("Category", value=p_source.get('category', 'General'))
            
            if st.button("💾 SAVE TO LIBRARY"):
                # Commit to DB
                final_meta = p_source['metadata']
                final_meta['title'] = p_title
                final_meta['project'] = p_project
                final_meta['category'] = p_category
                
                db_manager.save_transcript(p_source['url'], p_source['text'], metadata=final_meta)
                del st.session_state['preview_source'] # Clear preview
                st.success("Saved!")
                st.rerun()

    st.markdown("---")

    # --- 2. PROJECT EXPLORER (Tree View) ---
    st.markdown("#### 📂 MY PROJECTS")
    all_sources = db_manager.get_all_transcripts()
    if not all_sources:
        st.info("No projects yet.")
        selected_source = None
    else:
        # Group by Project -> Category
        projects = {}
        for s in all_sources:
            meta = s.get('metadata', {})
            proj = meta.get('project', 'Uncategorized')
            cat = meta.get('category', 'General')
            if proj not in projects: projects[proj] = {}
            if cat not in projects[proj]: projects[proj][cat] = []
            projects[proj][cat].append(s)
        
        # Render Tree
        selected_source = None
        # Use session state to track selection to avoid reset on rerun
        if 'selected_source_id' not in st.session_state: st.session_state['selected_source_id'] = None

        for proj_name, categories in projects.items():
            with st.expander(f"📁 {proj_name}", expanded=False):
                for cat_name, s_list in categories.items():
                    st.markdown(f"**{cat_name}**")
                    for s in s_list:
                        # Item Button
                        title = s.get('metadata', {}).get('title', s['url'][:15])
                        if st.button(f"📄 {title}", key=s['id']):
                            st.session_state['selected_source_id'] = s['id']
                            st.rerun()
        
        # Retrieve actual selected object
        if st.session_state['selected_source_id']:
            selected_source = next((s for s in all_sources if s['id'] == st.session_state['selected_source_id']), None)

    # --- 3. SKILL LAUNCHER ---
    st.markdown("---")
    st.markdown("#### 🧠 MY SKILLS")
    if not os.path.exists("skills"): os.makedirs("skills")
    skills = [f for f in os.listdir("skills") if f.endswith(".md")]
    
    if skills:
        selected_skill = st.selectbox("Select Skill", skills, label_visibility="collapsed")
        if st.button("🚀 LAUNCH SKILL"):
             st.session_state['active_skill'] = selected_skill
    else:
        st.caption("No skills created yet.")

# --- RIGHT: MAIN WORKSPACE (The "Wide Box") ---
with workspace_col: # Renamed from studio_col to match layout definition
    
    # DETERMINE ACTIVE CONTENT: Database Selection OR Live Preview
    active_source = selected_source
    if not active_source and 'preview_source' in st.session_state:
        active_source = st.session_state['preview_source']
        st.info("⚡ PREVIEW MODE - Save to Library to keep this.")
    
    # Alias back to selected_source so downstream logic works
    selected_source = active_source

    if selected_source:
        # --- 🎧 MASTER MEDIA PLAYER (Persistent Header) ---
        meta = selected_source.get('metadata', {})
        media_path = meta.get('video_path') or meta.get('audio_path')
        
        # Resolve to absolute path
        base_dir = "/Users/jaydengle/Transcribe-Reels"
        
        # Heuristic: If path missing/broken, search in 'reel' folder
        if not media_path or not os.path.exists(media_path):
            reel_dir = os.path.join(base_dir, 'reel')
            if os.path.exists(reel_dir):
                possible = [os.path.join(reel_dir, f) for f in os.listdir(reel_dir)
                            if f.endswith(('.mp4', '.mkv', '.webm', '.mov'))]
                if possible: media_path = possible[0]

        # Check for TTS audio podcast
        tts_path = os.path.join(base_dir, 'overview_audio.mp3')
        edu_path = os.path.join(base_dir, 'educational_podcast.mp3')
        has_tts = os.path.exists(tts_path)
        has_edu = os.path.exists(edu_path)
        
        with st.expander("🎧 MEDIA PANEL", expanded=True):
            panel_col1, panel_col2, panel_col3 = st.columns([1.2, 1, 1])
            
            with panel_col1:
                st.markdown("**🎥 Source Footage**")
                if media_path and os.path.exists(media_path):
                    st.video(media_path)
                else:
                    st.info("💡 Video cleared from storage.")
            
            with panel_col2:
                st.markdown("**🎧 Quick AI Podcast**")
                st.caption("Auto-generated 30s summary")
                if has_tts:
                    st.audio(tts_path)
                elif 'tts_audio_bytes' in st.session_state:
                    st.audio(st.session_state['tts_audio_bytes'], format='audio/mp3')
                else:
                    st.info("Appears after ⚡ Analyze.")
            
            with panel_col3:
                st.markdown("**🎓 Educational Podcast**")
                st.caption("2-speaker deep dive (~3 min)")
                if has_edu:
                    st.audio(edu_path)
                    if st.button("🔄 Regenerate", key="regen_edu"):
                        os.remove(edu_path)
                        st.rerun()
                else:
                    if openai_client and st.button("🎙️ Generate Educational Podcast", key="gen_edu", use_container_width=True):
                        with st.spinner("🎙️ Generating 3-min dual-speaker podcast..."):
                            t_text = selected_source.get('text', '')
                            a_text = selected_source.get('summary', '') or st.session_state.get('viral_analysis', '')
                            result = generate_educational_podcast(t_text, a_text, openai_client)
                            if result and not result.startswith("ERROR"):
                                st.success("🎓 Podcast Ready!")
                                st.rerun()
                            else:
                                st.error(f"Failed: {result}")
                    elif not openai_client:
                        st.info("OpenAI key needed.")

        # TABS interface for maximum space efficiency
        tab_transcript, tab_viral, tab_campaign, tab_visuals, tab_skill, tab_chat = st.tabs([
            "📝 TRANSCRIPT", 
            "🧬 VIRAL DNA", 
            "📢 CAMPAIGN", 
            "🎨 VISUALS",
            "🧠 SKILLS",
            "💬 AI CHAT"
        ])
        
        # --- TAB 1: TRANSCRIPT ---
        with tab_transcript:
            st.subheader("📄 Original Content")
            
            # Dynamic Hallucination Check for all items
            t_text = selected_source.get('text', '')
            JUNK_TERMS = ["SNOWY", "THANK YOU", "PAST IN THE GAMMA'S", "ANYWAY IM THE FAMOUS", "THE END", "TICK TOK"]
            is_risky = any(term in t_text.upper() for term in JUNK_TERMS) or len(t_text) < 20
            
            if is_risky:
                st.warning("⚠️ **Ambiguous Transcription**\nThe audio reflects background music or low speech. Please verify with the player below.")
            
            # A. Two columns: Text + Audio side by side
            txt_col, aud_col = st.columns([3, 2])
            
            with txt_col:
                edited_text = st.text_area("Edit Transcript", value=t_text, height=400)
            
            with aud_col:
                st.markdown("**🎧 Audio Transcript**")
                st.caption("Listen to the transcript read aloud")
                
                audio_transcript_path = "/Users/jaydengle/Transcribe-Reels/audio_transcript.mp3"
                
                if os.path.exists(audio_transcript_path):
                    st.audio(audio_transcript_path)
                    if st.button("🔄 Re-generate Audio", key="regen_audio_transcript"):
                        os.remove(audio_transcript_path)
                        st.rerun()
                else:
                    if openai_client and t_text:
                        if st.button("🎧 Generate Audio Transcript", key="gen_audio_transcript", use_container_width=True):
                            with st.spinner("🎧 Generating audio..."):
                                try:
                                    tts_resp = openai_client.audio.speech.create(
                                        model="tts-1",
                                        voice="alloy",
                                        input=t_text[:4096]  # TTS limit
                                    )
                                    tts_resp.stream_to_file(audio_transcript_path)
                                    st.success("🎧 Audio ready!")
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"TTS failed: {e}")
                    else:
                        st.info("Transcript needed first.")
        
        # --- TAB 2: VIRAL DNA (Wide Analysis) ---
        with tab_viral:
            st.subheader("🧬 Viral Structure Analysis")
            col1, col2 = st.columns([1, 3])
            
            with col1:
                st.markdown("**Actions**")
                
                # A. Text Analysis Trigger
                if st.button("🔍 Re-Analyze Structure", use_container_width=True):
                    if openai_client:
                        with st.spinner("Decoding Viral Patterns..."):
                            sys_prompt = "You are a Viral Content Analyst."
                            usr_prompt = f"""
                            Analyze this transcript and output a structured report with these Exact Headers:

                            ## 🪝 VIRAL HOOKS
                            (List 3 alternative viral hooks for this specific concept)

                            ## 🧠 WHY IT WENT VIRAL
                            (The psychological triggers: Curiosity, Fear Of Missing Out, Status, etc.)

                            ## 📐 CONTENT STRATEGY
                            (The structure used: e.g. Hook -> Story -> Value -> Call To Action)

                            {selected_source.get('text')}
                            """
                            resp = openai_client.chat.completions.create(model="gpt-4", messages=[{"role": "system", "content": sys_prompt}, {"role": "user", "content": usr_prompt}])
                            # Update Session State AND Database
                            new_analysis = resp.choices[0].message.content
                            st.session_state['viral_analysis'] = new_analysis
                            # Optional: Update DB here if we had an update function

                # A2. Mind Map Trigger
                if st.button("🕸️ Visual Mind Map (Flow)", use_container_width=True):
                     if openai_client:
                         with st.spinner("Mapping Viral Flow..."):
                             dot_prompt = f"""
                             Create a Graphviz DOT code to visualize the narrative flow of this transcript.
                             Nodes should be: Hook, Context, Problem, Solution, CTA.
                             Make it simple and linear.
                             TRANSCRIPT: {selected_source.get('text')}
                             OUTPUT ONLY THE DOT CODE. BEGIN with 'digraph'.
                             """
                             resp = openai_client.chat.completions.create(model="gpt-4", messages=[{"role": "system", "content": "You are a Graphviz Generator. Output DOT code only."}, {"role": "user", "content": dot_prompt}])
                             dot_code = resp.choices[0].message.content.replace("```dot", "").replace("```", "").strip()
                             st.session_state['viral_map'] = dot_code
                            
                # B. Visual Analysis Trigger
                if st.button("👁️ Deep Visual Analysis (GPT-4o)", use_container_width=True):
                    # ... (Existing Logic) ...
                    # Re-using logic to get video_path if needed
                    meta = selected_source.get('metadata', {})
                    video_path = meta.get('video_path')
                    if not video_path or not os.path.exists(video_path):
                         # Auto-fetch logic (re-using same pattern as before)
                         re_download = download_reel(selected_source.get('url'))
                         if re_download: video_path = re_download["video_path"]
                    
                    if video_path and os.path.exists(video_path) and openai_client:
                        with st.spinner("👁️ Decoding Visual Patterns (0s, 3s, 10s)..."):
                            try:
                                frames = analyze_visual_frames(video_path)
                                prompt_messages = [
                                    {
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": f"These are frames from a specific Reel (0s, 3s, 10s, End). Analyze the VISUAL HOOK. What happens in the first 3 seconds visually? How does the camera move? What text is overlayed? Explain why the VISUALS made this viral."},
                                            *map(lambda x: {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{x}"}}, frames)
                                        ],
                                    }
                                ]
                                resp = openai_client.chat.completions.create(model="gpt-4o", messages=prompt_messages, max_tokens=600)
                                st.session_state['visual_deep_dive'] = resp.choices[0].message.content
                            except Exception as e:
                                st.error(f"Visual Analysis Error: {e}")

                # B2. NEW: VISUAL TIMELINE
                if st.button("🎞️ Generate Frame-by-Frame Timeline", use_container_width=True):
                     meta = selected_source.get('metadata', {})
                     video_path = meta.get('video_path')
                     # Auto-fetch if missing
                     if not video_path or not os.path.exists(video_path):
                         re_download = download_reel(selected_source.get('url'))
                         if re_download: video_path = re_download["video_path"]

                     if video_path and os.path.exists(video_path) and openai_client:
                         with st.spinner("🎞️  Extracting & Analyzing Timeline (This takes ~30s)..."):
                             frames, stamps = generate_visual_timeline(video_path, interval=5)
                             
                             # Analyze with GPT-4o
                             prompt_content = [{"type": "text", "text": "Analyze this sequence of frames from a video (every 5 seconds). Describe exactly what is happening visually in each shot to explain the storytelling flow."}]
                             for i, f in enumerate(frames):
                                 prompt_content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{f}"}})
                                 prompt_content.append({"type": "text", "text": f"Frame at {stamps[i]}s"})
                             
                             resp = openai_client.chat.completions.create(model="gpt-4o", messages=[{"role": "user", "content": prompt_content}], max_tokens=1000)
                             
                             st.session_state['visual_timeline'] = {
                                 "analysis": resp.choices[0].message.content,
                                 "frames": frames,
                                 "timestamps": stamps
                             }

                st.markdown("---")
                
                # C. NEW: Save as Skill
                if st.button("🧠 Extract into Skill", use_container_width=True):
                     if openai_client and 'viral_analysis' in st.session_state:
                         with st.spinner("Compiling Skill..."):
                             skill_prompt = f"Create a reusable SKILL.md instruction file based on this analysis:\n{st.session_state['viral_analysis']}"
                             resp = openai_client.chat.completions.create(model="gpt-4", messages=[{"role": "user", "content": skill_prompt}])
                             skill_content = resp.choices[0].message.content
                             
                             # Save file
                             safe_title = "".join(x for x in selected_source.get('metadata', {}).get('title', 'skill') if x.isalnum())
                             with open(f"skills/{safe_title}.md", "w") as f:
                                 f.write(skill_content)
                             st.success(f"Skill Saved: {safe_title}")

                if st.button("📝 Generate Copycat Script", use_container_width=True):
                     if openai_client:
                         with st.spinner("Writing Script..."):
                             sys_prompt = "You are a Master Scriptwriter. Recreate this format for a new topic."
                             usr_prompt = f"Rewrite this video structure for a generic Personal Brand regarding 'Growth':\n{selected_source.get('text')}"
                             resp = openai_client.chat.completions.create(model="gpt-4", messages=[{"role": "system", "content": sys_prompt}, {"role": "user", "content": usr_prompt}])
                             st.session_state['recreation_script'] = resp.choices[0].message.content

            with col2:
                # The WIDE Display Area
                
                # 0. Mind Map Output (NEW)
                if 'viral_map' in st.session_state:
                    st.success("🕸️ Narrative Flow")
                    try:
                        st.graphviz_chart(st.session_state['viral_map'])
                    except Exception as e:
                        st.error(f"Graph Error: {e}")
                    
                    if 'viral_map' in st.session_state:
                         st.download_button("📥 Download Graph (DOT)", st.session_state['viral_map'], "viral_flow.dot")
                    
                    st.divider()

                # Check for existing analysis results
                if 'visual_timeline' in st.session_state:
                    st.success("🎞️ Visual Timeline Analysis")
                    data = st.session_state['visual_timeline']
                    st.markdown(data['analysis'])
                    
                    st.markdown("### 🎞️ Frame Sequence")
                    frames = data['frames']
                    stamps = data['timestamps']
                    
                    # Display in rows of 3
                    cols = st.columns(3)
                    for i, f in enumerate(frames):
                        with cols[i % 3]:
                            st.image(f"data:image/jpeg;base64,{f}", caption=f"{stamps[i]}s", use_container_width=True)
                    st.divider()

                if 'visual_deep_dive' in st.session_state:
                    st.success("👁️ Visual Decoder Output")
                    st.markdown(st.session_state['visual_deep_dive'])
                    st.divider()

                # 2. Text Analysis Output (Use Stored Summary if available, or session state)
                summary = selected_source.get('summary')
                if summary: 
                    st.session_state['viral_analysis'] = summary # Sync DB to Session
                
                if 'viral_analysis' in st.session_state:
                    st.info("✅ Text Structure Analysis")
                    st.markdown(st.session_state['viral_analysis'])
                elif not summary:
                     st.warning("⚠️ No Analysis Found. Click 'Re-Analyze Structure'.")
                
                # 3. Script Output
                if 'recreation_script' in st.session_state:
                    st.success("📝 Recreation Script")
                    st.markdown(st.session_state['recreation_script'])

        # --- TAB 3: CAMPAIGN LAUNCHER ---
        with tab_campaign:
             st.subheader("📢 Campaign Engineering")
             
             c_col1, c_col2 = st.columns([1, 2])
             with c_col1:
                 product_name = st.text_input("Product / Niche", placeholder="e.g. Hair Shampoo")
                 if st.button("🚀 Launch Campaign", use_container_width=True):
                    if not product_name:
                        st.warning("Enter a product name.")
                    elif openai_client:
                        with st.spinner(f"Engineering Campaign..."):
                             sys_prom = "You are a World-Class Marketing Strategist."
                             usr_prom = f"Take the VIRAL STRUCTURE of this transcript: '{selected_source.get('text')}' and create a 3-video mini-campaign for '{product_name}' (Hook, Value, Offer)."
                             resp = openai_client.chat.completions.create(model="gpt-4", messages=[{"role": "system", "content": sys_prom}, {"role": "user", "content": usr_prom}])
                             st.session_state['campaign_result'] = resp.choices[0].message.content
             
             with c_col2:
                 if 'campaign_result' in st.session_state:
                     st.markdown(st.session_state['campaign_result'])
                 else:
                     st.info("Enter a product and launch to see the generated strategy here.")

        # --- TAB 4: VISUAL STUDIO ---
        with tab_visuals:
            st.subheader("🎨 Visuals & Deck")
            v_col1, v_col2 = st.columns([1, 2])
            
            with v_col1:
                ppt_style = st.selectbox("Style", 
                    ["Cyberpunk Matrix", "Neon Synthwave", "Corporate Clean", "Dark Mode", "Golden Luxury"])
                
                if st.button("🖼 Generate Cover Art (OpenAI)", use_container_width=True):
                    with st.spinner("Dreaming up cover art (DALL-E 3)..."):
                         if openai_client:
                             try:
                                 art_prompt = f"Create a high-contrast YouTube Thumbnail for a video about: {selected_source.get('text')[:100]}. Style: {ppt_style}. No text."
                                 resp = openai_client.images.generate(model="dall-e-3", prompt=art_prompt, size="1024x1024", quality="standard", n=1)
                                 st.session_state['generated_image'] = resp.data[0].url
                             except Exception as e:
                                 st.error(f"Image Gen Error: {e}")

                if st.button("✨ Generate with Pollinations (Unlimited)", use_container_width=True):
                    with st.spinner("Pollinations (Flux) generating asset..."):
                        # Get a detailed prompt first
                        base_text = selected_source.get('text', '')[:200]
                        detailed_prompt = generate_image_desc_gemini(base_text)
                        
                        # Generate image via Pollinations
                        img_uri = generate_image_pollinations(detailed_prompt + f". Style: {ppt_style}.")
                        if img_uri:
                            st.session_state['generated_image'] = img_uri
                            st.success("Pollinations Asset Ready!")
                        else:
                            st.warning("Image Gen failed.")

                st.divider()
                st.markdown("### 🎥 Kling AI Storyboard")
                sb_engine = st.radio("Select Engine", ["Gemini 1.5 Flash (Free/Fast)", "GPT-4o Vision (Premium/Precise)"], horizontal=True)

                if st.button("🎥 Generate Kling AI Storyboard", use_container_width=True):
                    meta = selected_source.get('metadata', {})
                    v_path = meta.get('video_path')
                    
                    # Auto-find video if path broken
                    if not v_path or not os.path.exists(v_path):
                        reel_dir = "/Users/jaydengle/Transcribe-Reels/reel"
                        if os.path.exists(reel_dir):
                            possible = [os.path.join(reel_dir, f) for f in os.listdir(reel_dir)
                                        if f.endswith(('.mp4', '.mkv', '.webm', '.mov'))]
                            if possible: v_path = possible[0]

                    if not v_path or not os.path.exists(v_path):
                        st.error("❌ No video file found. Please RE-ANALYZE the URL to download the video.")
                    else:
                        with st.spinner("🎥 Analyzing video structure..."):
                            try:
                                result = []
                                if "Gemini" in sb_engine:
                                    # Use new Gemini Native function
                                    result = generate_kling_storyboard_gemini(v_path)
                                else:
                                    # Use legacy GPT-4o function (requires key)
                                    if not openai_client:
                                        st.error("❌ OpenAI API Key Missing.")
                                    else:
                                        result = generate_kling_storyboard(v_path, openai_client)

                                if isinstance(result, list):
                                    st.session_state['kling_storyboard'] = result
                                    st.success(f"🎥 Storyboard ready — {len(result)} shots!")
                                    st.rerun()
                                else:
                                    st.error(f"Analysis Failed: {result}")
                            except Exception as e:
                                st.error(f"Critical Error: {e}")

                # Display storyboard if generated
                if 'kling_storyboard' in st.session_state:
                    sb = st.session_state['kling_storyboard']
                    st.markdown("---")
                    st.markdown("**🎥 Your Kling AI Storyboard — Copy prompts directly into Kling Story Mode:**")
                    
                    all_prompts = ""
                    for shot in sb:
                        sb_c1, sb_c2 = st.columns([2, 1])
                        with sb_c1:
                            cut_badge = " ⚡ SCENE CUT" if shot.get('is_scene_cut') else ""
                            st.markdown(f"**Shot {shot['shot']}{cut_badge} — [{shot['timestamp']}s, ~{shot['duration']}s]**")
                            st.code(shot['prompt'], language=None)
                            all_prompts += f"\n### Shot {shot['shot']} [{shot['timestamp']}s]\n{shot['prompt']}\n"
                    
                    # Export as text file
                    st.download_button(
                        label="⬇️ Export All Prompts (.txt)",
                        data=all_prompts,
                        file_name="kling_storyboard.txt",
                        mime="text/plain",
                        key="dl_kling"
                    )
                
                st.divider()
                if st.button("📄 Generate How-To PDF (Frame-by-Frame)", use_container_width=True):
                    meta = selected_source.get('metadata', {})
                    v_path = meta.get('video_path')
                    
                    if not v_path or not os.path.exists(v_path):
                        reel_dir = "/Users/jaydengle/Transcribe-Reels/reel"
                        if os.path.exists(reel_dir):
                            possible = [os.path.join(reel_dir, f) for f in os.listdir(reel_dir)
                                        if f.endswith(('.mp4', '.mkv', '.webm', '.mov'))]
                            if possible: v_path = possible[0]
                    
                    if not openai_client:
                        st.error("❌ OpenAI API Key Missing. Please check your .env or sidebar settings.")
                    elif not v_path or not os.path.exists(v_path):
                        st.error("❌ No video file found. Please RE-ANALYZE the URL to download the video.")
                    else:
                        with st.spinner("📄 GPT-4o Vision reading every frame... (~60s)"):
                            try:
                                t_text = selected_source.get('text', '')
                                result = generate_how_to_pdf(v_path, t_text, openai_client)
                                if result and not result.startswith("ERROR"):
                                    with open(result, 'rb') as fp:
                                        st.download_button(
                                            label="⬇️ Download How-To PDF",
                                            data=fp.read(),
                                            file_name="how_to_guide.pdf",
                                            mime="application/pdf",
                                            key="dl_pdf"
                                        )
                                    st.success("📄 PDF Ready! Click above to download.")
                                else:
                                    st.error(f"PDF Analysis Failed: {result}")
                            except Exception as e:
                                st.error(f"Critical Error: {e}")

                if st.button("📊 Create Slide Deck (PPTX)", use_container_width=True):
                     if openai_client:
                         with st.spinner("Compiling Deck..."):
                            try:
                                from pptx import Presentation
                                from pptx.util import Inches, Pt
                                from pptx.dml.color import RGBColor
                                import io
                                import requests
                                import json

                                # 1. Generate Outline
                                full_context = f"Transcript:\n{selected_source.get('text')}\n\n"
                                instruction = "Create a powerful 7-slide presentation outline based on this text. Return a JSON list of objects only."
                                ppt_prompt = f"{instruction} \n\nCONTEXT:\n{full_context}"
                                response = openai_client.chat.completions.create(model="gpt-4", response_format={ "type": "json_object" }, messages=[{"role": "user", "content": ppt_prompt}])
                                slides_data = json.loads(response.choices[0].message.content)
                                
                                # 2. Build PPTX (Basic Implementation for Stability)
                                prs = Presentation()
                                prs.slide_width = Inches(16)
                                prs.slide_height = Inches(9)
                                blank_slide_layout = prs.slide_layouts[6] 
                                
                                # Content Slides
                                slides_list = slides_data["slides"] if "slides" in slides_data else slides_data.values()
                                for slide_info in list(slides_list)[:7]: 
                                    slide = prs.slides.add_slide(blank_slide_layout)
                                    # Title
                                    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(2))
                                    p = textbox.text_frame.add_paragraph()
                                    p.text = slide_info.get("title", "Insight").upper()
                                    p.font.size = Pt(40)
                                    
                                    # Content
                                    cbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
                                    p = cbox.text_frame.add_paragraph()
                                    p.text = str(slide_info.get("content", ""))

                                # Save
                                pptx_stream = io.BytesIO()
                                prs.save(pptx_stream)
                                pptx_stream.seek(0)
                                st.session_state['pptx_file'] = pptx_stream
                                st.success("Deck Ready!")

                            except Exception as e:
                                st.error(f"PPTX Error: {e}")

                if 'pptx_file' in st.session_state:
                    st.download_button(
                        label="📥 Download PPTX",
                        data=st.session_state['pptx_file'],
                        file_name="reel_insights.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )

            with v_col2:
                if 'generated_image' in st.session_state:
                     st.image(st.session_state['generated_image'], use_container_width=True)
                else:
                     st.info("Generated assets will appear here.")

        # --- TAB 5: SKILLS & TRAINING ---
        with tab_skill:
            st.subheader("🧠 Skill Studio")
            
            # ============================================================
            # SECTION 1: REVERSE ENGINEER ANY SCREENSHOT
            # ============================================================
            with st.expander("🔍 REVERSE ENGINEER ANY WORKFLOW", expanded=True):
                st.markdown("Upload **any screenshot** of a tool, workflow, or technique — the AI will identify it and tell you exactly how to replicate it.")
                re_img = st.file_uploader("Drop a screenshot here", type=["png","jpg","jpeg","webp"], key="re_uploader")
                
                if re_img:
                    up_col, res_col = st.columns([1, 2])
                    with up_col:
                        st.image(re_img, caption="Your Screenshot", use_container_width=True)
                    
                    with res_col:
                        if st.button("🔍 Reverse Engineer This", use_container_width=True, key="do_reverse"):
                            with st.spinner("🔍 GPT-4o Vision analyzing..."):
                                analysis = reverse_engineer_image(re_img.read(), openai_client)
                                st.session_state['reverse_analysis'] = analysis
                    
                    if 'reverse_analysis' in st.session_state:
                        analysis_md = st.session_state['reverse_analysis']
                        st.markdown(analysis_md)
                        
                        # Extract GitHub search terms and auto-search
                        if "GitHub Search Terms" in analysis_md:
                            search_block = analysis_md.split("GitHub Search Terms")[-1].strip().split('\n')
                            search_terms = [l.strip().strip('-').strip() for l in search_block if l.strip() and not l.startswith('#')]
                            
                            if search_terms:
                                st.divider()
                                st.markdown("**📦 Auto-Discovered GitHub Resources:**")
                                first_term = search_terms[0] if search_terms else ""
                                if first_term:
                                    with st.spinner(f"🔎 Searching GitHub for: {first_term}"):
                                        repos = search_github_for_skill(first_term)
                                    
                                    if repos:
                                        for repo in repos:
                                            r_col1, r_col2 = st.columns([3, 1])
                                            with r_col1:
                                                st.markdown(f"**[{repo['name']}]({repo['url']})** ⭐ {repo['stars']:,}")
                                                st.caption(repo['description'])
                                            with r_col2:
                                                skill_content = f"# {repo['name']}\n\n{repo['description']}\n\n**Stars:** {repo['stars']:,}\n\n**GitHub:** {repo['url']}\n\n**Clone:** `git clone {repo['clone']}`\n\n---\n\n## How This Relates\n\n{analysis_md[:500]}"
                                                if st.button("💾 Save Skill", key=f"save_{repo['name'].replace('/','_')}"):
                                                    safe_name = repo['name'].replace('/', '_')
                                                    if not os.path.exists("/Users/jaydengle/Transcribe-Reels/skills"):
                                                        os.makedirs("/Users/jaydengle/Transcribe-Reels/skills")
                                                    with open(f"/Users/jaydengle/Transcribe-Reels/skills/{safe_name}.md", "w") as f:
                                                        f.write(skill_content)
                                                    st.success("Saved!")
                                    else:
                                        st.info("No repos found — try a different search term.")
                        
                        # Save full analysis as skill
                        st.divider()
                        if st.button("💾 Save Full Analysis as Skill Module", key="save_reverse_skill", use_container_width=True):
                            if not os.path.exists("/Users/jaydengle/Transcribe-Reels/skills"):
                                os.makedirs("/Users/jaydengle/Transcribe-Reels/skills")
                            fname = f"/Users/jaydengle/Transcribe-Reels/skills/reverse_engineer_{int(__import__('time').time())}.md"
                            with open(fname, "w") as f:
                                f.write(f"# Reverse Engineered Skill\n\n{analysis_md}")
                            st.success(f"Saved to skills library!")

            st.divider()
            
            # ============================================================
            # SECTION 2: EXTRACT SKILLS FROM CURRENT VIDEO
            # ============================================================
            st.markdown("**🧠 Extract Skills from Current Video (Gemini)**")
            s_col1, s_col2 = st.columns([1, 2])
            with s_col1:
                st.info("Uses Gemini frame analysis to identify technical skills demonstrated in the loaded video.")
                if st.button("🧠 Extract Exact Skills", use_container_width=True):
                    if 'visual_timeline' in st.session_state:
                        frames = st.session_state['visual_timeline']['frames']
                        with st.spinner("Gemini analyzing skills..."):
                            skills_text = extract_skills_gemini(frames)
                            st.session_state['extracted_skills'] = skills_text
                    else:
                        st.warning("Please generate a Visual Timeline in Tab 2 first so I have frames to analyze!")
            
            with s_col2:
                if 'extracted_skills' in st.session_state:
                    st.success("Extracted Skills & Techniques")
                    st.markdown(st.session_state['extracted_skills'])
                    if st.button("💾 Save as Skill Module", use_container_width=True):
                        safe_title = "".join(x for x in selected_source.get('metadata', {}).get('title', 'skill') if x.isalnum())
                        if not os.path.exists("/Users/jaydengle/Transcribe-Reels/skills"):
                            os.makedirs("/Users/jaydengle/Transcribe-Reels/skills")
                        with open(f"/Users/jaydengle/Transcribe-Reels/skills/{safe_title}_skills.md", "w") as f:
                            f.write(st.session_state['extracted_skills'])
                        st.success(f"Skill Module Saved!")
                elif 'active_skill' in st.session_state:
                    skill_file = st.session_state['active_skill']
                    st.info(f"Loaded Skill: {skill_file}")
                    try:
                        with open(f"/Users/jaydengle/Transcribe-Reels/skills/{skill_file}", "r") as f:
                            st.markdown(f.read())
                    except: st.error("Skill file missing.")
                else:
                    st.info("Select a skill from the Left Sidebar or Extract from current video.")

        # --- TAB 6: AI CHAT (Restored) ---
        with tab_chat:
            st.subheader("💬 Neural Chat")
            # Initialize chat history
            if "messages" not in st.session_state:
                st.session_state.messages = []

            # Display chat messages from history on app rerun
            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            # React to user input
            if prompt := st.chat_input("Ask about this video..."):
                # Display user message in chat message container
                st.chat_message("user").markdown(prompt)
                # Add user message to chat history
                st.session_state.messages.append({"role": "user", "content": prompt})

                # Assistant response
                if openai_client:
                    with st.chat_message("assistant"):
                        message_placeholder = st.empty()
                        full_response = ""
                        
                        # Context: The Transcript + Analysis
                        context = f"Context:\nTranscript: {selected_source.get('text')}\nAnalysis: {selected_source.get('summary', '')}"
                        
                        stream = openai_client.chat.completions.create(
                            model="gpt-4",
                            messages=[
                                {"role": "system", "content": f"You are a helpful assistant analyzing this video context:\n{context}"},
                                *[{"role": m["role"], "content": m["content"]} for m in st.session_state.messages]
                            ],
                            stream=True,
                        )
                        
                        for chunk in stream:
                            if chunk.choices[0].delta.content:
                                full_response += chunk.choices[0].delta.content
                                message_placeholder.markdown(full_response + "▌")
                        message_placeholder.markdown(full_response)
                        
                    st.session_state.messages.append({"role": "assistant", "content": full_response})
    else:
        st.info("👈 Please select or add a source to begin.")
