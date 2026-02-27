import os
import sys
import json
import time
import shutil
from datetime import datetime

import yt_dlp
from moviepy import VideoFileClip
from openai import OpenAI
from supabase import create_client

try:
    import whisper
except Exception:
    whisper = None

try:
    from scenedetect import open_video, SceneManager
    from scenedetect.detectors import ContentDetector
except Exception:
    open_video = None
    SceneManager = None
    ContentDetector = None

try:
    from graphviz import Source
except Exception:
    Source = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from fpdf import FPDF
except Exception:
    FPDF = None

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
EXPORT_DIR = os.path.join(BASE_DIR, "export")


def ensure_dir(path):
    if not os.path.exists(path):
        os.makedirs(path, exist_ok=True)


def log_step(label, ok, detail=None):
    status = "✅" if ok else "❌"
    print(f"{status} {label}")
    if detail:
        print(f"   {detail}")


def download_video(url, out_dir):
    ensure_dir(out_dir)
    ydl_opts = {
        "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
        "outtmpl": os.path.join(out_dir, "%(title)s.%(ext)s"),
        "quiet": True,
        "no_warnings": True,
        "user_agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "referer": "https://www.instagram.com/",
        "cookiesfrombrowser": ("chrome",),
    }

    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        info = ydl.extract_info(url, download=True)
        title = info.get("title", "video")
        description = info.get("description", "")
        uploader = info.get("uploader", "Unknown")

    video_exts = [".mp4", ".mkv", ".webm", ".mov"]
    downloaded = [f for f in os.listdir(out_dir) if any(f.endswith(ext) for ext in video_exts)]
    if not downloaded:
        raise RuntimeError("No video file downloaded.")

    video_path = os.path.join(out_dir, downloaded[0])
    return {
        "title": title,
        "caption": description,
        "owner": uploader,
        "video_path": video_path,
    }


def extract_audio(video_path, out_dir):
    audio_path = os.path.join(out_dir, "audio.wav")
    clip = VideoFileClip(video_path)
    clip.audio.write_audiofile(audio_path, codec="pcm_s16le", ffmpeg_params=["-ac", "1", "-ar", "16000"], logger=None)
    clip.close()
    return audio_path


def transcribe_audio(audio_path, model_size="base"):
    if whisper is None:
        raise RuntimeError("openai-whisper is not installed.")
    model = whisper.load_model(model_size)
    result = model.transcribe(audio_path)
    return {
        "text": result.get("text", "").strip(),
        "language": result.get("language"),
    }


def openai_chat(client, prompt, model="gpt-4o"):
    resp = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=900,
    )
    return resp.choices[0].message.content


def generate_mindmap_dot(client, transcript):
    prompt = (
        "You are a Graphviz Generator. Output DOT code only. "
        "Create a concise mind map of the key ideas from this transcript:\n\n"
        f"{transcript[:3500]}"
    )
    raw = openai_chat(client, prompt, model="gpt-4o")
    return sanitize_dot(raw)


def write_text(path, content):
    with open(path, "w") as f:
        f.write(content.strip() + "\n")


def sanitize_dot(dot_text):
    cleaned = dot_text.strip()
    if "```" in cleaned:
        parts = cleaned.split("```")
        if len(parts) >= 2:
            cleaned = parts[1].strip()
    if cleaned.startswith("dot"):
        cleaned = cleaned.replace("dot", "", 1).strip()
    return cleaned


def render_graphviz(dot, out_dir):
    if Source is None:
        raise RuntimeError("graphviz package not installed.")
    graph = Source(dot)
    png_path = os.path.join(out_dir, "mindmap.png")
    graph.render(filename=png_path, format="png", cleanup=True)
    if not png_path.endswith(".png"):
        png_path += ".png"
    return png_path


def generate_pdf(summary, out_dir):
    if FPDF is None:
        raise RuntimeError("fpdf2 not installed.")
    pdf = FPDF()
    pdf.add_page()
    font_path = "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"
    if os.path.exists(font_path):
        pdf.add_font("ArialUnicode", "", font_path, uni=True)
        pdf.set_font("ArialUnicode", size=12)
    else:
        pdf.set_font("Helvetica", size=12)
    pdf.multi_cell(0, 8, "Video Summary")
    pdf.ln(2)
    pdf.multi_cell(0, 8, summary)
    pdf_path = os.path.join(out_dir, "summary.pdf")
    pdf.output(pdf_path)
    return pdf_path


def generate_pptx(summary, out_dir):
    if Presentation is None:
        raise RuntimeError("python-pptx not installed.")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Video Summary"
    slide.placeholders[1].text = summary[:1500]
    pptx_path = os.path.join(out_dir, "summary.pptx")
    prs.save(pptx_path)
    return pptx_path


def generate_tts(client, text, out_dir):
    tts_path = os.path.join(out_dir, "summary.mp3")
    resp = client.audio.speech.create(
        model="tts-1",
        voice="alloy",
        input=text[:2000],
    )
    resp.stream_to_file(tts_path)
    return tts_path


def detect_scenes(video_path, out_dir):
    if open_video is None:
        raise RuntimeError("scenedetect not available.")
    video = open_video(video_path)
    scene_mgr = SceneManager()
    scene_mgr.add_detector(ContentDetector(threshold=27.0))
    scene_mgr.detect_scenes(video, show_progress=False)
    scenes = scene_mgr.get_scene_list()
    scene_data = [
        {
            "start": s[0].get_seconds(),
            "end": s[1].get_seconds(),
        }
        for s in scenes
    ]
    out_path = os.path.join(out_dir, "scenes.json")
    with open(out_path, "w") as f:
        json.dump(scene_data, f, indent=2)
    return out_path


def save_to_supabase(url, text, summary, metadata):
    supa_url = os.getenv("SUPABASE_URL")
    supa_key = os.getenv("SUPABASE_KEY")
    if not supa_url or not supa_key:
        raise RuntimeError("Supabase env vars not set.")
    client = create_client(supa_url, supa_key)
    data = {
        "url": url,
        "text": text,
        "summary": summary,
        "metadata": metadata,
        "created_at": datetime.utcnow().isoformat(),
    }
    client.table("transcripts").upsert(data, on_conflict="url").execute()


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python full_run_test.py <video_url>")
        sys.exit(1)

    url = sys.argv[1]
    run_id = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    out_dir = os.path.join(EXPORT_DIR, f"run_{run_id}")
    ensure_dir(out_dir)

    openai_key = os.getenv("OPENAI_API_KEY")
    if not openai_key:
        print("Missing OPENAI_API_KEY.")
        sys.exit(1)

    client = OpenAI(api_key=openai_key)

    # Step 1: Download
    try:
        video_meta = download_video(url, out_dir)
        log_step("Download video", True, video_meta.get("video_path"))
    except Exception as e:
        log_step("Download video", False, str(e))
        sys.exit(1)

    # Step 2: Extract audio
    try:
        audio_path = extract_audio(video_meta["video_path"], out_dir)
        log_step("Extract audio", True, audio_path)
    except Exception as e:
        log_step("Extract audio", False, str(e))
        sys.exit(1)

    # Step 3: Transcribe
    try:
        transcription = transcribe_audio(audio_path, model_size="base")
        transcript_text = transcription.get("text", "")
        log_step("Transcribe audio", True, f"{len(transcript_text)} chars")
    except Exception as e:
        log_step("Transcribe audio", False, str(e))
        sys.exit(1)

    # Step 4: Summary + Viral analysis
    try:
        summary_prompt = "Summarize the transcript in 8 bullet points.\n\n" + transcript_text[:3500]
        summary = openai_chat(client, summary_prompt, model="gpt-4o")
        log_step("Generate summary", True)
    except Exception as e:
        log_step("Generate summary", False, str(e))
        summary = None

    # Step 4b: Detailed report
    try:
        report_prompt = (
            "Create a detailed report with sections: Overview, Key Points, "
            "Actionable Tips, Risks/Mistakes, and Final Takeaways.\n\n"
            + transcript_text[:3500]
        )
        report = openai_chat(client, report_prompt, model="gpt-4o")
        report_path = os.path.join(out_dir, "report.txt")
        write_text(report_path, report)
        log_step("Generate report", True, report_path)
    except Exception as e:
        log_step("Generate report", False, str(e))

    # Step 5: Mind map
    try:
        dot = generate_mindmap_dot(client, transcript_text)
        dot_path = os.path.join(out_dir, "mindmap.dot")
        with open(dot_path, "w") as f:
            f.write(dot)
        log_step("Generate mind map DOT", True, dot_path)
        if Source:
            png_path = render_graphviz(dot, out_dir)
            log_step("Render mind map PNG", True, png_path)
    except Exception as e:
        log_step("Generate mind map", False, str(e))

    # Step 5b: Viral analysis + hooks
    try:
        viral_prompt = (
            "Analyze why this video could go viral. Provide: Hook, "
            "Retention triggers, Emotional drivers, and Rewatch moments.\n\n"
            + transcript_text[:3500]
        )
        viral_analysis = openai_chat(client, viral_prompt, model="gpt-4o")
        viral_path = os.path.join(out_dir, "viral_analysis.txt")
        write_text(viral_path, viral_analysis)
        log_step("Generate viral analysis", True, viral_path)

        hooks_prompt = (
            "List 8 alternative hook lines for this video, each under 12 words.\n\n"
            + transcript_text[:2000]
        )
        hooks = openai_chat(client, hooks_prompt, model="gpt-4o")
        hooks_path = os.path.join(out_dir, "hooks.txt")
        write_text(hooks_path, hooks)
        log_step("Generate hooks", True, hooks_path)
    except Exception as e:
        log_step("Generate viral analysis/hooks", False, str(e))

    # Step 6: Scene detection
    try:
        scenes_path = detect_scenes(video_meta["video_path"], out_dir)
        log_step("Detect scenes", True, scenes_path)
    except Exception as e:
        log_step("Detect scenes", False, str(e))

    # Step 7: PDF and PPTX
    if summary:
        try:
            pdf_path = generate_pdf(summary, out_dir)
            log_step("Generate PDF", True, pdf_path)
        except Exception as e:
            log_step("Generate PDF", False, str(e))

        try:
            pptx_path = generate_pptx(summary, out_dir)
            log_step("Generate PPTX", True, pptx_path)
        except Exception as e:
            log_step("Generate PPTX", False, str(e))

        try:
            tts_path = generate_tts(client, summary, out_dir)
            log_step("Generate TTS", True, tts_path)
        except Exception as e:
            log_step("Generate TTS", False, str(e))

    # Step 7b: Podcast script
    try:
        podcast_prompt = (
            "Write a 3-minute two-speaker podcast script based on this transcript. "
            "Use labels HOST: and EXPERT:.\n\n"
            + transcript_text[:3500]
        )
        podcast = openai_chat(client, podcast_prompt, model="gpt-4o")
        podcast_path = os.path.join(out_dir, "podcast_script.txt")
        write_text(podcast_path, podcast)
        log_step("Generate podcast script", True, podcast_path)
    except Exception as e:
        log_step("Generate podcast script", False, str(e))

    # Step 7c: Skills extraction + copycat plan
    try:
        skills_prompt = (
            "Extract a concise skill list from the transcript and then provide a "
            "copycat plan: steps to reproduce the skill in a new context.\n\n"
            + transcript_text[:3500]
        )
        skills = openai_chat(client, skills_prompt, model="gpt-4o")
        skills_path = os.path.join(out_dir, "skills_copycat.txt")
        write_text(skills_path, skills)
        log_step("Generate skills/copycat", True, skills_path)
    except Exception as e:
        log_step("Generate skills/copycat", False, str(e))

    # Step 7d: Strategy + campaign + recreate prompt
    try:
        strategy_prompt = (
            "Create a content strategy for a 7-day plan based on this transcript. "
            "Include audience, angles, and distribution tips.\n\n"
            + transcript_text[:3500]
        )
        strategy = openai_chat(client, strategy_prompt, model="gpt-4o")
        strategy_path = os.path.join(out_dir, "strategy.txt")
        write_text(strategy_path, strategy)
        log_step("Generate strategy", True, strategy_path)

        campaign_prompt = (
            "Create a 3-video mini-campaign (Hook, Value, Offer) based on this transcript.\n\n"
            + transcript_text[:3500]
        )
        campaign = openai_chat(client, campaign_prompt, model="gpt-4o")
        campaign_path = os.path.join(out_dir, "campaign.txt")
        write_text(campaign_path, campaign)
        log_step("Generate campaign", True, campaign_path)

        recreate_prompt = (
            "Create a shot-by-shot recreation plan (5-7 shots) and a single master prompt "
            "to recreate this video.\n\n"
            + transcript_text[:3500]
        )
        recreate = openai_chat(client, recreate_prompt, model="gpt-4o")
        recreate_path = os.path.join(out_dir, "recreate.txt")
        write_text(recreate_path, recreate)
        log_step("Generate recreate plan", True, recreate_path)
    except Exception as e:
        log_step("Generate strategy/campaign/recreate", False, str(e))

    # Step 8: Save to Supabase
    try:
        metadata = {
            "title": video_meta.get("title"),
            "owner": video_meta.get("owner"),
            "caption": video_meta.get("caption"),
        }
        save_to_supabase(url, transcript_text, summary, metadata)
        log_step("Save to Supabase", True)
    except Exception as e:
        log_step("Save to Supabase", False, str(e))

    # Copy outputs to a stable path
    latest_dir = os.path.join(EXPORT_DIR, "latest")
    if os.path.exists(latest_dir):
        shutil.rmtree(latest_dir)
    shutil.copytree(out_dir, latest_dir)

    print(f"\nRun complete. Outputs in: {out_dir}")
    print(f"Latest outputs in: {latest_dir}")
